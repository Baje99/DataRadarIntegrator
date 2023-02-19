from pathlib import Path
from tkinter import Tk, Canvas, Entry, Button, PhotoImage
from tkinter.filedialog import askopenfilename
import pandas as pd
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor 
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN

OUTPUT_PATH = Path(__file__).parent
ASSETS_PATH = OUTPUT_PATH / Path(r"#PATH of downloaded asstes ( assets\frame0)")

def relative_to_assets(path: str) -> Path:
    return ASSETS_PATH / Path(path)

def CreateInterface():
    global entry_1, canvas
    window = Tk()
    window.geometry("862x519")
    window.configure(bg = "#0F1827")
    canvas = Canvas(
        window,
        bg = "#0F1827",
        height = 519,
        width = 862,
        bd = 0,
        highlightthickness = 0,
        relief = "ridge"
    )

    canvas.place(x = 0, y = 0)
    canvas.create_rectangle(
        430.9999999999999,
        7.105427357601002e-15,
        861.9999999999999,
        519.0,
        fill="#FCFCFC",
        outline="")

    button_image_1 = PhotoImage(
        file=relative_to_assets("button_1.png"))
    button_1 = Button(
        image=button_image_1,
        borderwidth=0,
        highlightthickness=0,
        command=IntegrateData,
        relief="flat"
    )
    button_1.place(
        x=558.9999999999999,
        y=391.0,
        width=176.0,
        height=51.0
    )

    canvas.create_rectangle(
        477.9999999999999,
        213.0,
        822.9999999999999,
        274.0,
        fill="#F1F5FF",
        outline="")

    button_image_2 = PhotoImage(
        file=relative_to_assets("button_2.png"))
    button_2 = Button(
        image=button_image_2,
        borderwidth=0,
        highlightthickness=0,
        command=UploadPowerpoint,
        relief="flat"
    )
    button_2.place(
        x=782.9999999999999,
        y=233.0,
        width=24.0,
        height=22.0
    )

    canvas.create_text(
        495.9999999999999,
        230.0,
        anchor="nw",
        text="Upload PowerPoint",
        fill="#505485",
        font=("Roboto Bold", 24 * -1)
    )

    canvas.create_rectangle(
        477.9999999999999,
        127.0,
        822.9999999999999,
        188.0,
        fill="#F1F5FF",
        outline="")

    button_image_3 = PhotoImage(
        file=relative_to_assets("button_3.png"))
    button_3 = Button(
        image=button_image_3,
        borderwidth=0,
        highlightthickness=0,
        command=UploadExcel,
        relief="flat"
    )
    button_3.place(
        x=782.9999999999999,
        y=147.0,
        width=24.0,
        height=22.0
    )

    canvas.create_text(
        496.9999999999999,
        144.0,
        anchor="nw",
        text="Upload Excel",
        fill="#505485",
        font=("Roboto Bold", 24 * -1)
    )

    canvas.create_rectangle(
        477.9999999999999,
        296.0,
        822.9999999999999,
        343.0,
        fill="#F1F5FF",
        outline="")

    canvas.create_text(
        491.9999999999999,
        306.0,
        anchor="nw",
        text="Enter Slide Number:",
        fill="#505485",
        font=("Roboto Bold", 24 * -1)
    )

    canvas.create_text(
        491.9999999999999,
        45.00000000000001,
        anchor="nw",
        text="Radar Data Integrator",
        fill="#515486",
        font=("Roboto Bold", 32 * -1)
    )

    entry_image_1 = PhotoImage(
        file=relative_to_assets("entry_1.png"))
    entry_bg_1 = canvas.create_image(
        792.4999999999999,
        320.0,
        image=entry_image_1
    )
    entry_1 = Entry(
        bd=0,
        bg="#515486",
        fg="#000716",
        highlightthickness=0
    )
    entry_1.place(
        x=780.9999999999999,
        y=299.0,
        width=23.0,
        height=40.0
    )

    image_image_1 = PhotoImage(
        file=relative_to_assets("image_1.png"))
    image_1 = canvas.create_image(
        215.0000000000001,
        259.0,
        image=image_image_1
    )
    window.resizable(False, False)
    window.mainloop()

def UploadExcel():
    global dct
    dct = {"Time": [], "Course": [], "Altitude": [], "Speed": [], "Azimuth": []}
    f_types = [('Excel files',"*.xlsx"),("Csv files", "*.csv")]
    file = askopenfilename(initialdir=os.getcwd(), 
                       filetypes=f_types,
                       title='Upload ODAS')
    file_name, file_extension = os.path.splitext(file)
    speccolumns = [x for x in range(1,6)]
    if file_extension == ".csv":
        df=pd.read_csv(file, usecols = speccolumns)
    else:
        df=pd.read_excel(file, usecols = speccolumns)
    for colname, colval in df.items():
        min = colval.values[0]
        max = colval.values[0]
        for item in colval.values:
            if item > max:
                max = item
            elif item < min:
                min = item
        if colname == "Time":
            min = min.strftime("%H%M")
            max = max.strftime("%H%M")
        dct[colname].append(str(min))
        dct[colname].append(str(max))

def UploadPowerpoint():
    global pptfile
    f_types = [('PPTX',"*.pptx"),("PPTM", "*.pptm"), ("PPT", "*.ppt")]
    pptfile = askopenfilename(initialdir=os.getcwd(), 
                       filetypes=f_types,
                       title='Upload POWERPOINT')

def IntegrateData():
    try :
        input = int(entry_1.get())
    except ValueError:
        print("Please enter a number")
    prs = Presentation(pptfile)
    slide = prs.slides[input]
    shapes = slide.shapes
    width = Inches(5.5)
    height = Inches(1.0)
    left  = Inches(4.0)
    top = Inches(3.0)
    shape = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()                                
    shape.fill.fore_color.rgb = RGBColor(150, 100, 100)
    pg = shape.text_frame.paragraphs[0]
    pg.alignment = PP_ALIGN.LEFT
    run = pg.add_run()   	
    run.text = "T: " + dct["Time"][0]  + " - " + dct["Time"][1] + ";Course:" + dct["Course"][0] + " - " + dct["Course"][1] + ";Speed:" + dct["Speed"][0] + " - " + dct["Speed"][1] + "km/h" + ";Azimuth:" + dct["Azimuth"][0] + " - " + dct["Azimuth"][1] + ";Altitude:" + dct["Altitude"][0] + " - " + dct["Altitude"][1] + "m"
    font = run.font      
    font.bold = True
    font.name = 'Arial Narrow'
    font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)
    font.size = Pt(11)  
    prs.save("Example2.pptx")
    canvas.create_text(
        580.9999999999999,
        445.00000000000001,
        anchor="nw",
        text="Operation successful!",
        fill="#515486",
        font=("Roboto Bold", 16 * -1)
    )

if __name__ == "__main__":
    CreateInterface()


